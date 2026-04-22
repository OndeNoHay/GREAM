"""
pptx_graem MCP server — genera presentaciones PowerPoint (.pptx) con estilos ATEXIS.

Herramientas:
  create_presentation — crea un .pptx desde slides JSON estructuradas
  list_templates      — lista plantillas .potx disponibles
"""

import json
import pathlib
from typing import Optional, Union

from fastmcp import FastMCP
from mcp.types import ToolAnnotations

mcp = FastMCP(name="pptx_graem")

_PROJECT_ROOT = pathlib.Path(__file__).parent.parent.parent.resolve()
_OUTPUT_DIR = _PROJECT_ROOT / "output"
_TEMPLATES_DIR = _PROJECT_ROOT / "templates"

# ATEXIS brand colors
_DARK_BLUE = (31, 56, 100)   # #1F3864
_LIGHT_BLUE = (46, 117, 182)  # #2E75B6
_ORANGE = (237, 125, 49)     # #ED7D31

# PowerPoint default slide layout indices
_LAYOUT_TITLE = 0    # Title Slide (placeholder 0=title, 1=subtitle)
_LAYOUT_CONTENT = 1  # Title and Content (placeholder 0=title, 1=body)
_LAYOUT_BLANK = 6    # Blank


def _ensure_output_dir() -> pathlib.Path:
    _OUTPUT_DIR.mkdir(parents=True, exist_ok=True)
    return _OUTPUT_DIR


def _set_text_color(paragraph, rgb: tuple) -> None:
    from pptx.dml.color import RGBColor
    for run in paragraph.runs:
        run.font.color.rgb = RGBColor(*rgb)


def _build_presentation(title: str, slides: list, template_path: Optional[pathlib.Path] = None):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    prs = Presentation(str(template_path)) if (template_path and template_path.exists()) else Presentation()

    layout_map = {
        "title": _LAYOUT_TITLE,
        "content": _LAYOUT_CONTENT,
        "blank": _LAYOUT_BLANK,
    }
    n_layouts = len(prs.slide_layouts)

    for slide_data in slides:
        layout_key = slide_data.get("layout", "content")
        layout_idx = min(layout_map.get(layout_key, _LAYOUT_CONTENT), n_layouts - 1)
        layout = prs.slide_layouts[layout_idx]
        slide = prs.slides.add_slide(layout)
        slide_title = slide_data.get("title", "")

        if layout_key == "title":
            for ph in slide.placeholders:
                idx = ph.placeholder_format.idx
                if idx == 0:
                    ph.text = slide_title
                    _set_text_color(ph.text_frame.paragraphs[0], _DARK_BLUE)
                elif idx == 1:
                    ph.text = slide_data.get("subtitle", "")

        elif layout_key == "content":
            body_lines = slide_data.get("body", [])
            if isinstance(body_lines, str):
                body_lines = [body_lines]
            for ph in slide.placeholders:
                idx = ph.placeholder_format.idx
                if idx == 0:
                    ph.text = slide_title
                    _set_text_color(ph.text_frame.paragraphs[0], _DARK_BLUE)
                elif idx == 1:
                    tf = ph.text_frame
                    tf.clear()
                    for i, line in enumerate(body_lines):
                        if i == 0:
                            tf.paragraphs[0].text = str(line)
                        else:
                            tf.add_paragraph().text = str(line)

        elif layout_key == "blank" and slide_title:
            txBox = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(1))
            tf = txBox.text_frame
            tf.text = slide_title
            p = tf.paragraphs[0]
            if p.runs:
                run = p.runs[0]
                run.font.size = Pt(28)
                run.font.bold = True
                run.font.color.rgb = RGBColor(*_DARK_BLUE)

    return prs


@mcp.tool(annotations=ToolAnnotations(destructiveHint=True))
def create_presentation(
    output_filename: str,
    title: str,
    slides_json: str,
    template: str = "",
) -> str:
    """
    Crea una presentación PowerPoint (.pptx) con estilos ATEXIS.

    slides_json es un array JSON con diapositivas:
      {"layout": "title",   "title": "...", "subtitle": "..."}
      {"layout": "content", "title": "...", "body": ["Punto 1", "Punto 2"]}
      {"layout": "blank",   "title": "..."}

    Si el primer slide no es de tipo "title", se añade automáticamente una
    portada con el título de la presentación.

    template: nombre de fichero .potx en templates/ (vacío = estilos por defecto).
    Devuelve la ruta absoluta del fichero creado.
    """
    try:
        slides = json.loads(slides_json) if slides_json.strip() else []
    except json.JSONDecodeError as e:
        return f"Error: slides_json inválido — {e}"

    if not output_filename.endswith(".pptx"):
        output_filename += ".pptx"

    safe_name = pathlib.Path(output_filename).name
    if safe_name != output_filename:
        return "Error: output_filename debe ser solo nombre de fichero, sin directorios"

    if not slides or slides[0].get("layout") != "title":
        slides = [{"layout": "title", "title": title, "subtitle": "ATEXIS Group"}] + slides

    template_path = None
    if template:
        candidate = (_TEMPLATES_DIR / template).resolve()
        if str(candidate).startswith(str(_TEMPLATES_DIR)) and candidate.exists():
            template_path = candidate

    try:
        prs = _build_presentation(title, slides, template_path)
        out_path = _ensure_output_dir() / safe_name
        prs.save(str(out_path))
        return str(out_path)
    except Exception as e:
        return f"Error al crear presentación: {e}"


@mcp.tool(annotations=ToolAnnotations(readOnlyHint=True, idempotentHint=True))
def list_templates() -> str:
    """Lista las plantillas PowerPoint (.potx) disponibles en el directorio templates/."""
    if not _TEMPLATES_DIR.exists():
        return json.dumps({"templates": [], "count": 0})
    templates = [
        {"name": f.name, "size_bytes": f.stat().st_size}
        for f in sorted(_TEMPLATES_DIR.iterdir())
        if f.is_file() and f.suffix.lower() == ".potx"
    ]
    return json.dumps({"templates": templates, "count": len(templates)})


if __name__ == "__main__":
    mcp.run()
