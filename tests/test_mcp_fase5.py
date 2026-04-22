"""
Tests de Fase 5 — word_graem + pptx_graem MCP servers.

Cubre:
  - word_graem: create_document, create_s1000d_changelog, list_templates
  - pptx_graem: create_presentation, list_templates
  - Integración: arranque real vía MCPClientManager
"""

import json
import pathlib

import pytest


# ---------------------------------------------------------------------------
# Fixture: redirige _OUTPUT_DIR y _TEMPLATES_DIR a tmp_path
# ---------------------------------------------------------------------------

@pytest.fixture(autouse=True)
def patch_dirs(tmp_path, monkeypatch):
    import mcp_servers.word_graem.server as wmod
    import mcp_servers.pptx_graem.server as pmod

    out = tmp_path / "output"
    tpl = tmp_path / "templates"
    out.mkdir()
    tpl.mkdir()

    monkeypatch.setattr(wmod, "_OUTPUT_DIR", out)
    monkeypatch.setattr(wmod, "_TEMPLATES_DIR", tpl)
    monkeypatch.setattr(pmod, "_OUTPUT_DIR", out)
    monkeypatch.setattr(pmod, "_TEMPLATES_DIR", tpl)


# ---------------------------------------------------------------------------
# word_graem — create_document
# ---------------------------------------------------------------------------

class TestWordCreateDocument:
    def test_creates_docx_file(self):
        from mcp_servers.word_graem.server import create_document
        content = json.dumps([
            {"type": "heading", "level": 1, "text": "Introduction"},
            {"type": "paragraph", "text": "Test paragraph."},
        ])
        result = create_document("test_doc.docx", "Test Document", content)
        assert not result.startswith("Error"), result
        assert pathlib.Path(result).exists()

    def test_extension_auto_appended(self):
        from mcp_servers.word_graem.server import create_document
        result = create_document("no_ext", "Title", "[]")
        assert not result.startswith("Error"), result
        assert result.endswith(".docx")

    def test_empty_content(self):
        from mcp_servers.word_graem.server import create_document
        result = create_document("empty.docx", "Title", "[]")
        assert not result.startswith("Error"), result
        assert pathlib.Path(result).exists()

    def test_invalid_json_returns_error(self):
        from mcp_servers.word_graem.server import create_document
        result = create_document("x.docx", "Title", "{bad json")
        assert result.startswith("Error")
        assert "inválido" in result

    def test_path_traversal_blocked(self):
        from mcp_servers.word_graem.server import create_document
        result = create_document("../evil.docx", "Title", "[]")
        assert result.startswith("Error")

    def test_nested_path_blocked(self):
        from mcp_servers.word_graem.server import create_document
        result = create_document("subdir/file.docx", "Title", "[]")
        assert result.startswith("Error")

    def test_table_content(self):
        from mcp_servers.word_graem.server import create_document
        from docx import Document
        content = json.dumps([{
            "type": "table",
            "headers": ["DMC", "Title", "Issue"],
            "rows": [
                ["DMC-ATEST-A-32-00-00-00A-040A-D", "Hydraulic", "001"],
                ["DMC-ATEST-A-32-10-00-00A-040A-D", "Landing Gear", "001"],
            ],
        }])
        result = create_document("table_test.docx", "Table Doc", content)
        doc = Document(result)
        assert len(doc.tables) == 1
        assert doc.tables[0].rows[0].cells[0].text == "DMC"
        assert len(doc.tables[0].rows) == 3  # 1 header + 2 data rows

    def test_heading_levels_and_paragraph(self):
        from mcp_servers.word_graem.server import create_document
        from docx import Document
        content = json.dumps([
            {"type": "heading", "level": 1, "text": "H1"},
            {"type": "heading", "level": 2, "text": "H2"},
            {"type": "heading", "level": 3, "text": "H3"},
            {"type": "paragraph", "text": "Body text"},
        ])
        result = create_document("structure.docx", "Structured", content)
        doc = Document(result)
        texts = [p.text for p in doc.paragraphs]
        assert "H1" in texts
        assert "H2" in texts
        assert "Body text" in texts

    def test_pagebreak_accepted(self):
        from mcp_servers.word_graem.server import create_document
        content = json.dumps([
            {"type": "paragraph", "text": "Before break"},
            {"type": "pagebreak"},
            {"type": "paragraph", "text": "After break"},
        ])
        result = create_document("pagebreak.docx", "PB Test", content)
        assert not result.startswith("Error"), result


# ---------------------------------------------------------------------------
# word_graem — create_s1000d_changelog
# ---------------------------------------------------------------------------

class TestWordCreateChangelog:
    def test_creates_changelog(self):
        from mcp_servers.word_graem.server import create_s1000d_changelog
        from docx import Document
        entries = json.dumps([
            {"issue": "001", "date": "2024-01-15", "description": "Initial release",
             "author": "JJO", "reason": "New document"},
            {"issue": "002", "date": "2024-02-01", "description": "Updated hydraulic section",
             "author": "JJO", "reason": "Technical update"},
        ])
        result = create_s1000d_changelog("changelog.docx", entries)
        assert not result.startswith("Error"), result
        doc = Document(result)
        assert len(doc.tables) == 1
        assert len(doc.tables[0].rows) == 3  # header + 2 entries

    def test_empty_entries(self):
        from mcp_servers.word_graem.server import create_s1000d_changelog
        result = create_s1000d_changelog("empty_cl.docx", "[]")
        assert not result.startswith("Error"), result

    def test_invalid_json_returns_error(self):
        from mcp_servers.word_graem.server import create_s1000d_changelog
        result = create_s1000d_changelog("cl.docx", "not valid json")
        assert result.startswith("Error")

    def test_custom_document_title(self):
        from mcp_servers.word_graem.server import create_s1000d_changelog
        from docx import Document
        result = create_s1000d_changelog("cl2.docx", "[]", document_title="S1000D Change Record")
        doc = Document(result)
        texts = [p.text for p in doc.paragraphs]
        assert any("S1000D Change Record" in t for t in texts)

    def test_path_traversal_blocked(self):
        from mcp_servers.word_graem.server import create_s1000d_changelog
        result = create_s1000d_changelog("../../evil.docx", "[]")
        assert result.startswith("Error")

    def test_header_columns_present(self):
        from mcp_servers.word_graem.server import create_s1000d_changelog
        from docx import Document
        result = create_s1000d_changelog("cl3.docx", "[]")
        doc = Document(result)
        assert len(doc.tables) == 1
        headers = [c.text for c in doc.tables[0].rows[0].cells]
        assert "Issue" in headers
        assert "Description" in headers
        assert "Author" in headers


# ---------------------------------------------------------------------------
# word_graem — list_templates
# ---------------------------------------------------------------------------

class TestWordListTemplates:
    def test_empty_when_no_dotx(self):
        from mcp_servers.word_graem.server import list_templates
        data = json.loads(list_templates())
        assert data["count"] == 0
        assert data["templates"] == []

    def test_lists_only_dotx(self, tmp_path):
        import mcp_servers.word_graem.server as wmod
        (wmod._TEMPLATES_DIR / "atexis.dotx").write_bytes(b"fake")
        (wmod._TEMPLATES_DIR / "readme.docx").write_bytes(b"fake")  # must be excluded
        from mcp_servers.word_graem.server import list_templates
        data = json.loads(list_templates())
        assert data["count"] == 1
        assert data["templates"][0]["name"] == "atexis.dotx"


# ---------------------------------------------------------------------------
# pptx_graem — create_presentation
# ---------------------------------------------------------------------------

class TestPptxCreatePresentation:
    def test_creates_pptx_file(self):
        from mcp_servers.pptx_graem.server import create_presentation
        slides = json.dumps([
            {"layout": "title", "title": "Main Title", "subtitle": "Sub"},
            {"layout": "content", "title": "Slide 2", "body": ["Point A", "Point B"]},
        ])
        result = create_presentation("test.pptx", "Test Pres", slides)
        assert not result.startswith("Error"), result
        assert pathlib.Path(result).exists()

    def test_extension_auto_appended(self):
        from mcp_servers.pptx_graem.server import create_presentation
        result = create_presentation("no_ext", "Title", "[]")
        assert not result.startswith("Error"), result
        assert result.endswith(".pptx")

    def test_auto_cover_slide_added(self):
        from mcp_servers.pptx_graem.server import create_presentation
        from pptx import Presentation
        slides = json.dumps([
            {"layout": "content", "title": "Content Slide", "body": ["item"]},
        ])
        result = create_presentation("auto_cover.pptx", "My Pres", slides)
        prs = Presentation(result)
        assert len(prs.slides) == 2  # auto cover + content slide

    def test_existing_title_slide_not_duplicated(self):
        from mcp_servers.pptx_graem.server import create_presentation
        from pptx import Presentation
        slides = json.dumps([
            {"layout": "title", "title": "Cover", "subtitle": "Sub"},
        ])
        result = create_presentation("no_dupe.pptx", "Cover", slides)
        prs = Presentation(result)
        assert len(prs.slides) == 1

    def test_empty_slides_creates_cover(self):
        from mcp_servers.pptx_graem.server import create_presentation
        from pptx import Presentation
        result = create_presentation("only_cover.pptx", "ATEXIS Demo", "[]")
        prs = Presentation(result)
        assert len(prs.slides) == 1

    def test_body_as_string_accepted(self):
        from mcp_servers.pptx_graem.server import create_presentation
        slides = json.dumps([
            {"layout": "content", "title": "Single body", "body": "Just a string"},
        ])
        result = create_presentation("str_body.pptx", "Title", slides)
        assert not result.startswith("Error"), result

    def test_blank_slide_with_title(self):
        from mcp_servers.pptx_graem.server import create_presentation
        slides = json.dumps([
            {"layout": "title", "title": "Cover"},
            {"layout": "blank", "title": "Appendix A"},
        ])
        result = create_presentation("blank_test.pptx", "Test", slides)
        assert not result.startswith("Error"), result
        assert pathlib.Path(result).exists()

    def test_invalid_json_returns_error(self):
        from mcp_servers.pptx_graem.server import create_presentation
        result = create_presentation("x.pptx", "Title", "{bad}")
        assert result.startswith("Error")

    def test_path_traversal_blocked(self):
        from mcp_servers.pptx_graem.server import create_presentation
        result = create_presentation("../evil.pptx", "Title", "[]")
        assert result.startswith("Error")

    def test_multi_slide_deck(self):
        from mcp_servers.pptx_graem.server import create_presentation
        from pptx import Presentation
        slides = json.dumps([
            {"layout": "title", "title": "GRAEM Demo", "subtitle": "ATEXIS Group"},
            {"layout": "content", "title": "Architecture", "body": ["MCP Layer", "FastAPI", "Kùzu"]},
            {"layout": "content", "title": "Results", "body": ["99% accuracy", "< 2s latency"]},
            {"layout": "blank", "title": "Q&A"},
        ])
        result = create_presentation("deck.pptx", "GRAEM Demo", slides)
        prs = Presentation(result)
        assert len(prs.slides) == 4


# ---------------------------------------------------------------------------
# pptx_graem — list_templates
# ---------------------------------------------------------------------------

class TestPptxListTemplates:
    def test_empty_when_no_potx(self):
        from mcp_servers.pptx_graem.server import list_templates
        data = json.loads(list_templates())
        assert data["count"] == 0
        assert data["templates"] == []

    def test_lists_only_potx(self, tmp_path):
        import mcp_servers.pptx_graem.server as pmod
        (pmod._TEMPLATES_DIR / "atexis.potx").write_bytes(b"fake")
        (pmod._TEMPLATES_DIR / "other.pptx").write_bytes(b"fake")  # must be excluded
        from mcp_servers.pptx_graem.server import list_templates
        data = json.loads(list_templates())
        assert data["count"] == 1
        assert data["templates"][0]["name"] == "atexis.potx"


# ---------------------------------------------------------------------------
# Integración: servidores arrancan vía MCPClientManager
# ---------------------------------------------------------------------------

@pytest.mark.integration
@pytest.mark.asyncio
async def test_word_graem_server_starts_and_lists_tools():
    from app.models.agents import MCPServerConfig
    from app.services.mcp_client_manager import MCPClientManager

    manager = MCPClientManager.__new__(MCPClientManager)
    manager._sessions = {}
    manager._contexts = {}
    manager._configs = {}

    config = MCPServerConfig(
        name="word_graem_test",
        type="stdio",
        command="python",
        args=["-m", "mcp_servers.word_graem.server"],
        enabled=True,
        timeout_seconds=30,
    )
    try:
        ok = await manager.start_server(config)
        assert ok, "word_graem server failed to start"
        tools = await manager.list_tools("word_graem_test")
        tool_names = {t.name for t in tools}
        assert "create_document" in tool_names
        assert "create_s1000d_changelog" in tool_names
        assert "list_templates" in tool_names
    finally:
        await manager.stop_all()


@pytest.mark.integration
@pytest.mark.asyncio
async def test_pptx_graem_server_starts_and_lists_tools():
    from app.models.agents import MCPServerConfig
    from app.services.mcp_client_manager import MCPClientManager

    manager = MCPClientManager.__new__(MCPClientManager)
    manager._sessions = {}
    manager._contexts = {}
    manager._configs = {}

    config = MCPServerConfig(
        name="pptx_graem_test",
        type="stdio",
        command="python",
        args=["-m", "mcp_servers.pptx_graem.server"],
        enabled=True,
        timeout_seconds=30,
    )
    try:
        ok = await manager.start_server(config)
        assert ok, "pptx_graem server failed to start"
        tools = await manager.list_tools("pptx_graem_test")
        tool_names = {t.name for t in tools}
        assert "create_presentation" in tool_names
        assert "list_templates" in tool_names
    finally:
        await manager.stop_all()
